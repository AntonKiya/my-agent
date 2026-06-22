import asyncio
from collections.abc import Iterable
from dataclasses import dataclass
from pathlib import Path
from typing import Protocol

from agent_service.media import MediaAsset

SUPPORTED_DOCUMENT_CONTENT_TYPES = frozenset(
    {
        "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        "application/vnd.openxmlformats-officedocument.presentationml.presentation",
        "application/pdf",
        "text/markdown",
        "text/plain",
        "text/x-markdown",
    }
)
SUPPORTED_DOCUMENT_SUFFIXES = frozenset({".docx", ".md", ".pdf", ".pptx", ".txt"})


class DocumentReadError(RuntimeError):
    def __init__(
        self,
        message: str,
        *,
        error_code: str | None = None,
    ) -> None:
        super().__init__(message)
        self.error_code = error_code


@dataclass(frozen=True, slots=True)
class DocumentContent:
    media_id: str
    filename: str | None
    content_type: str | None
    content: str
    truncated: bool
    size_bytes: int
    char_count: int


class DocumentReader(Protocol):
    def supports(self, asset: MediaAsset) -> bool:
        """Return whether this reader can extract content from the asset."""
        ...

    async def read(self, *, asset: MediaAsset, max_chars: int) -> DocumentContent:
        """Extract text content from the asset."""
        ...


@dataclass(frozen=True, slots=True)
class PlainTextDocumentReader:
    def supports(self, asset: MediaAsset) -> bool:
        filename = document_filename(asset)
        if not is_supported_document_payload(filename=filename, content_type=asset.content_type):
            return False
        suffix = _filename_suffix(filename)
        return suffix == ".txt" or (
            suffix is None
            and _normalized_content_type(asset.content_type) == "text/plain"
        )

    async def read(self, *, asset: MediaAsset, max_chars: int) -> DocumentContent:
        return await _read_text_asset(asset=asset, max_chars=max_chars)


@dataclass(frozen=True, slots=True)
class MarkdownDocumentReader:
    def supports(self, asset: MediaAsset) -> bool:
        filename = document_filename(asset)
        if not is_supported_document_payload(filename=filename, content_type=asset.content_type):
            return False
        suffix = _filename_suffix(filename)
        return suffix == ".md" or (
            suffix is None
            and _normalized_content_type(asset.content_type) in {"text/markdown", "text/x-markdown"}
        )

    async def read(self, *, asset: MediaAsset, max_chars: int) -> DocumentContent:
        return await _read_text_asset(asset=asset, max_chars=max_chars)


@dataclass(frozen=True, slots=True)
class DocxDocumentReader:
    def supports(self, asset: MediaAsset) -> bool:
        filename = document_filename(asset)
        if not is_supported_document_payload(filename=filename, content_type=asset.content_type):
            return False
        suffix = _filename_suffix(filename)
        return suffix == ".docx" or (
            suffix is None
            and _normalized_content_type(asset.content_type)
            == "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
        )

    async def read(self, *, asset: MediaAsset, max_chars: int) -> DocumentContent:
        return await asyncio.to_thread(_read_docx_asset, asset=asset, max_chars=max_chars)


@dataclass(frozen=True, slots=True)
class PdfDocumentReader:
    def supports(self, asset: MediaAsset) -> bool:
        filename = document_filename(asset)
        if not is_supported_document_payload(filename=filename, content_type=asset.content_type):
            return False
        suffix = _filename_suffix(filename)
        return suffix == ".pdf" or (
            suffix is None and _normalized_content_type(asset.content_type) == "application/pdf"
        )

    async def read(self, *, asset: MediaAsset, max_chars: int) -> DocumentContent:
        return await asyncio.to_thread(_read_pdf_asset, asset=asset, max_chars=max_chars)


@dataclass(frozen=True, slots=True)
class PptxDocumentReader:
    def supports(self, asset: MediaAsset) -> bool:
        filename = document_filename(asset)
        if not is_supported_document_payload(filename=filename, content_type=asset.content_type):
            return False
        suffix = _filename_suffix(filename)
        return suffix == ".pptx" or (
            suffix is None
            and _normalized_content_type(asset.content_type)
            == "application/vnd.openxmlformats-officedocument.presentationml.presentation"
        )

    async def read(self, *, asset: MediaAsset, max_chars: int) -> DocumentContent:
        return await asyncio.to_thread(_read_pptx_asset, asset=asset, max_chars=max_chars)


@dataclass(frozen=True, slots=True)
class DocumentReaderRegistry:
    readers: tuple[DocumentReader, ...] = (
        PlainTextDocumentReader(),
        MarkdownDocumentReader(),
        DocxDocumentReader(),
        PdfDocumentReader(),
        PptxDocumentReader(),
    )

    async def read(self, *, asset: MediaAsset, max_chars: int) -> DocumentContent:
        for reader in self.readers:
            if reader.supports(asset):
                return await reader.read(asset=asset, max_chars=max_chars)
        raise DocumentReadError(
            "Document format is not supported",
            error_code="unsupported_file_type",
        )


def is_supported_document_payload(*, filename: str | None, content_type: str | None) -> bool:
    suffix = _filename_suffix(filename)
    normalized_content_type = _normalized_content_type(content_type)
    if suffix is not None and suffix not in SUPPORTED_DOCUMENT_SUFFIXES:
        return False
    if normalized_content_type is not None:
        return normalized_content_type in SUPPORTED_DOCUMENT_CONTENT_TYPES
    return suffix in SUPPORTED_DOCUMENT_SUFFIXES


def document_filename(asset: MediaAsset) -> str | None:
    filename = asset.metadata.get("original_filename")
    if isinstance(filename, str) and filename:
        return filename
    candidate = Path(asset.storage_key).name
    return candidate or None


async def _read_text_asset(*, asset: MediaAsset, max_chars: int) -> DocumentContent:
    path = Path(asset.storage_key)
    try:
        content_bytes = await asyncio.to_thread(path.read_bytes)
    except OSError as exc:
        raise DocumentReadError(
            "Document file could not be read",
            error_code="file_read_failed",
        ) from exc

    content = _decode_text(content_bytes)
    truncated = len(content) > max_chars
    if truncated:
        content = content[:max_chars]

    return DocumentContent(
        media_id=asset.media_id,
        filename=document_filename(asset),
        content_type=asset.content_type,
        content=content,
        truncated=truncated,
        size_bytes=asset.size_bytes,
        char_count=len(content),
    )


def _read_docx_asset(*, asset: MediaAsset, max_chars: int) -> DocumentContent:
    path = Path(asset.storage_key)
    try:
        from docx import Document

        document = Document(str(path))
    except Exception as exc:
        raise DocumentReadError(
            "DOCX document could not be read",
            error_code="file_read_failed",
        ) from exc

    content = _extract_docx_content(document)
    truncated = len(content) > max_chars
    if truncated:
        content = content[:max_chars]

    return DocumentContent(
        media_id=asset.media_id,
        filename=document_filename(asset),
        content_type=asset.content_type,
        content=content,
        truncated=truncated,
        size_bytes=asset.size_bytes,
        char_count=len(content),
    )


def _read_pdf_asset(*, asset: MediaAsset, max_chars: int) -> DocumentContent:
    path = Path(asset.storage_key)
    try:
        from pypdf import PdfReader

        reader = PdfReader(str(path))
        pages = tuple(reader.pages)
        page_texts = []
        for page_number, page in enumerate(pages, start=1):
            text = (page.extract_text() or "").strip()
            if text:
                page_texts.append(f"[Page {page_number}]\n{text}")
    except Exception as exc:
        raise DocumentReadError(
            "PDF document could not be read",
            error_code="file_read_failed",
        ) from exc

    content = "\n\n".join(page_texts)
    if not content.strip():
        raise DocumentReadError(
            "PDF document does not contain extractable text",
            error_code="empty_file_content",
        )

    truncated = len(content) > max_chars
    if truncated:
        content = content[:max_chars]

    return DocumentContent(
        media_id=asset.media_id,
        filename=document_filename(asset),
        content_type=asset.content_type,
        content=content,
        truncated=truncated,
        size_bytes=asset.size_bytes,
        char_count=len(content),
    )


def _read_pptx_asset(*, asset: MediaAsset, max_chars: int) -> DocumentContent:
    path = Path(asset.storage_key)
    try:
        from pptx import Presentation

        presentation = Presentation(str(path))
        content = _extract_pptx_content(presentation)
    except Exception as exc:
        raise DocumentReadError(
            "PPTX document could not be read",
            error_code="file_read_failed",
        ) from exc

    if not content.strip():
        raise DocumentReadError(
            "PPTX document does not contain extractable text",
            error_code="empty_file_content",
        )

    truncated = len(content) > max_chars
    if truncated:
        content = content[:max_chars]

    return DocumentContent(
        media_id=asset.media_id,
        filename=document_filename(asset),
        content_type=asset.content_type,
        content=content,
        truncated=truncated,
        size_bytes=asset.size_bytes,
        char_count=len(content),
    )


def _extract_docx_content(document: object) -> str:
    blocks: list[str] = []
    paragraphs = getattr(document, "paragraphs", ())
    for paragraph in paragraphs:
        text = getattr(paragraph, "text", "").strip()
        if text:
            blocks.append(text)

    tables = getattr(document, "tables", ())
    for table_index, table in enumerate(tables, start=1):
        rows = []
        for row in getattr(table, "rows", ()):
            cells = [
                _normalize_table_cell_text(getattr(cell, "text", ""))
                for cell in getattr(row, "cells", ())
            ]
            if any(cells):
                rows.append("\t".join(cells))
        if rows:
            blocks.append(f"[Table {table_index}]\n" + "\n".join(rows))

    return "\n\n".join(blocks)


def _extract_pptx_content(presentation: object) -> str:
    blocks: list[str] = []
    slides = getattr(presentation, "slides", ())
    for slide_index, slide in enumerate(slides, start=1):
        slide_blocks, _ = _extract_pptx_shapes(getattr(slide, "shapes", ()))
        if slide_blocks:
            blocks.append(f"[Slide {slide_index}]\n" + "\n\n".join(slide_blocks))
    return "\n\n".join(blocks)


def _extract_pptx_shapes(
    shapes: Iterable[object],
    *,
    table_start: int = 1,
) -> tuple[list[str], int]:
    blocks: list[str] = []
    table_index = table_start
    for shape in shapes:
        child_shapes = getattr(shape, "shapes", None)
        if child_shapes is not None:
            child_blocks, table_index = _extract_pptx_shapes(
                child_shapes,
                table_start=table_index,
            )
            blocks.extend(child_blocks)

        if getattr(shape, "has_table", False):
            table_text = _extract_pptx_table(getattr(shape, "table", None))
            if table_text:
                blocks.append(f"[Table {table_index}]\n{table_text}")
                table_index += 1
            continue

        text = _extract_pptx_text(shape)
        if text:
            blocks.append(text)
    return blocks, table_index


def _extract_pptx_text(shape: object) -> str:
    if not getattr(shape, "has_text_frame", False):
        return ""
    text_frame = getattr(shape, "text_frame", None)
    lines = []
    for paragraph in getattr(text_frame, "paragraphs", ()):
        text = getattr(paragraph, "text", "").strip()
        if text:
            level = getattr(paragraph, "level", 0)
            indent = "  " * level if isinstance(level, int) else ""
            lines.append(f"{indent}{text}")
    return "\n".join(lines)


def _extract_pptx_table(table: object) -> str:
    rows = []
    for row in getattr(table, "rows", ()):
        cells = [
            _normalize_table_cell_text(getattr(cell, "text", ""))
            for cell in getattr(row, "cells", ())
        ]
        if any(cells):
            rows.append("\t".join(cells))
    return "\n".join(rows)


def _normalize_table_cell_text(text: str) -> str:
    return " ".join(part.strip() for part in text.splitlines() if part.strip())


def _decode_text(content: bytes) -> str:
    for encoding in ("utf-8-sig", "utf-8", "cp1251"):
        try:
            return content.decode(encoding)
        except UnicodeDecodeError:
            continue
    return content.decode("utf-8", errors="replace")


def _filename_suffix(filename: str | None) -> str | None:
    if filename is None:
        return None
    suffix = Path(filename).suffix.lower()
    return suffix or None


def _normalized_content_type(content_type: str | None) -> str | None:
    if content_type is None:
        return None
    normalized = content_type.split(";", 1)[0].strip().lower()
    return normalized or None
