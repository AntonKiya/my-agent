from dataclasses import dataclass, field
from pathlib import Path
from typing import Any, cast
from uuid import UUID, uuid4

from pydantic_ai import RunContext
from pydantic_ai.models.test import TestModel
from pydantic_ai.usage import RunUsage

from agent_service.config import AppSettings
from agent_service.document_reading.toolsets import build_file_reading_toolsets
from agent_service.media import MediaAsset, MediaAssetType


@dataclass(slots=True)
class FakeMediaAssetStore:
    assets: dict[tuple[str, UUID, UUID], MediaAsset] = field(default_factory=dict)

    async def create(self, *, asset: MediaAsset) -> MediaAsset:
        self.assets[(asset.media_id, asset.user_id, asset.conversation_id)] = asset
        return asset

    async def get(
        self,
        *,
        media_id: str,
        user_id: UUID,
        conversation_id: UUID,
    ) -> MediaAsset | None:
        return self.assets.get((media_id, user_id, conversation_id))


@dataclass(slots=True)
class FakeRunContext:
    deps: dict[str, Any]


async def test_file_reading_toolset_includes_usage_policy_instructions() -> None:
    toolsets = build_file_reading_toolsets(
        AppSettings(environment="test"),
        media_asset_store=FakeMediaAssetStore(),
    )
    instructions = await cast(Any, toolsets[0]).get_instructions(
        RunContext(deps={}, model=TestModel(), usage=RunUsage())
    )

    assert instructions is not None
    assert len(instructions) == 1
    content = instructions[0].content
    assert "previously attached file/document" in content
    assert "call readFileContent before answering" in content
    assert "Do not infer file contents from filename" in content


async def test_read_file_content_reads_owned_txt_file(tmp_path: Path) -> None:
    user_id = uuid4()
    conversation_id = uuid4()
    path = tmp_path / "notes.txt"
    path.write_text("hello from file", encoding="utf-8")
    store = FakeMediaAssetStore()
    await store.create(
        asset=MediaAsset(
            media_id="owned-file",
            user_id=user_id,
            conversation_id=conversation_id,
            media_type=MediaAssetType.DOCUMENT,
            storage_key=str(path),
            content_type="text/plain",
            size_bytes=path.stat().st_size,
            source_channel="telegram",
            metadata={"original_filename": "notes.txt"},
        )
    )
    toolsets = build_file_reading_toolsets(
        AppSettings(environment="test"),
        media_asset_store=store,
    )
    tool = cast(Any, toolsets[0]).tools["readFileContent"]
    ctx = cast(
        Any,
        FakeRunContext(deps={"user_id": user_id, "conversation_id": conversation_id}),
    )

    result = await tool.function(ctx, media_ids=["owned-file"])

    assert result["success"] is True
    assert result["data"]["files"] == [
        {
            "media_id": "owned-file",
            "filename": "notes.txt",
            "content_type": "text/plain",
            "content": "hello from file",
            "truncated": False,
            "size_bytes": path.stat().st_size,
            "char_count": len("hello from file"),
        }
    ]


async def test_read_file_content_reads_owned_docx_file(tmp_path: Path) -> None:
    user_id = uuid4()
    conversation_id = uuid4()
    path = tmp_path / "report.docx"
    _write_docx(path)
    store = FakeMediaAssetStore()
    await store.create(
        asset=MediaAsset(
            media_id="owned-docx",
            user_id=user_id,
            conversation_id=conversation_id,
            media_type=MediaAssetType.DOCUMENT,
            storage_key=str(path),
            content_type=(
                "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
            ),
            size_bytes=path.stat().st_size,
            source_channel="telegram",
            metadata={"original_filename": "report.docx"},
        )
    )
    toolsets = build_file_reading_toolsets(
        AppSettings(environment="test"),
        media_asset_store=store,
    )
    tool = cast(Any, toolsets[0]).tools["readFileContent"]
    ctx = cast(
        Any,
        FakeRunContext(deps={"user_id": user_id, "conversation_id": conversation_id}),
    )

    result = await tool.function(ctx, media_ids=["owned-docx"])

    assert result["success"] is True
    file = result["data"]["files"][0]
    assert file["media_id"] == "owned-docx"
    assert file["filename"] == "report.docx"
    assert file["content"] == "Project brief\n\n[Table 1]\nName\tValue\nCount\t42"
    assert file["truncated"] is False


async def test_read_file_content_reads_owned_pdf_file(tmp_path: Path) -> None:
    user_id = uuid4()
    conversation_id = uuid4()
    path = tmp_path / "report.pdf"
    _write_pdf(path, text="PDF project brief")
    store = FakeMediaAssetStore()
    await store.create(
        asset=MediaAsset(
            media_id="owned-pdf",
            user_id=user_id,
            conversation_id=conversation_id,
            media_type=MediaAssetType.DOCUMENT,
            storage_key=str(path),
            content_type="application/pdf",
            size_bytes=path.stat().st_size,
            source_channel="telegram",
            metadata={"original_filename": "report.pdf"},
        )
    )
    toolsets = build_file_reading_toolsets(
        AppSettings(environment="test"),
        media_asset_store=store,
    )
    tool = cast(Any, toolsets[0]).tools["readFileContent"]
    ctx = cast(
        Any,
        FakeRunContext(deps={"user_id": user_id, "conversation_id": conversation_id}),
    )

    result = await tool.function(ctx, media_ids=["owned-pdf"])

    assert result["success"] is True
    file = result["data"]["files"][0]
    assert file["media_id"] == "owned-pdf"
    assert file["filename"] == "report.pdf"
    assert file["content"] == "[Page 1]\nPDF project brief"
    assert file["truncated"] is False


async def test_read_file_content_reads_owned_pptx_file(tmp_path: Path) -> None:
    user_id = uuid4()
    conversation_id = uuid4()
    content_type = "application/vnd.openxmlformats-officedocument.presentationml.presentation"
    path = tmp_path / "deck.pptx"
    _write_pptx(path, include_content=True)
    store = FakeMediaAssetStore()
    await store.create(
        asset=MediaAsset(
            media_id="owned-pptx",
            user_id=user_id,
            conversation_id=conversation_id,
            media_type=MediaAssetType.DOCUMENT,
            storage_key=str(path),
            content_type=content_type,
            size_bytes=path.stat().st_size,
            source_channel="telegram",
            metadata={"original_filename": "deck.pptx"},
        )
    )
    toolsets = build_file_reading_toolsets(
        AppSettings(environment="test"),
        media_asset_store=store,
    )
    tool = cast(Any, toolsets[0]).tools["readFileContent"]
    ctx = cast(
        Any,
        FakeRunContext(deps={"user_id": user_id, "conversation_id": conversation_id}),
    )

    result = await tool.function(ctx, media_ids=["owned-pptx"])

    assert result["success"] is True
    file = result["data"]["files"][0]
    assert file["media_id"] == "owned-pptx"
    assert file["filename"] == "deck.pptx"
    assert file["content"] == (
        "[Slide 1]\n"
        "Project review\n\n"
        "Main goals\n  Increase retention\n\n"
        "[Table 1]\nMetric\tValue\nRetention\t42%"
    )
    assert file["truncated"] is False


async def test_read_file_content_checks_ownership_and_media_type(tmp_path: Path) -> None:
    user_id = uuid4()
    conversation_id = uuid4()
    other_user_id = uuid4()
    path = tmp_path / "notes.md"
    path.write_text("# Owned", encoding="utf-8")
    store = FakeMediaAssetStore()
    await store.create(
        asset=MediaAsset(
            media_id="other-file",
            user_id=other_user_id,
            conversation_id=conversation_id,
            media_type=MediaAssetType.DOCUMENT,
            storage_key=str(path),
            content_type="text/markdown",
            source_channel="telegram",
            metadata={"original_filename": "notes.md"},
        )
    )
    await store.create(
        asset=MediaAsset(
            media_id="image-id",
            user_id=user_id,
            conversation_id=conversation_id,
            media_type=MediaAssetType.IMAGE,
            storage_key=str(tmp_path / "image.jpg"),
            content_type="image/jpeg",
            source_channel="telegram",
        )
    )
    toolsets = build_file_reading_toolsets(
        AppSettings(environment="test"),
        media_asset_store=store,
    )
    tool = cast(Any, toolsets[0]).tools["readFileContent"]
    ctx = cast(
        Any,
        FakeRunContext(deps={"user_id": user_id, "conversation_id": conversation_id}),
    )

    denied_other = await tool.function(ctx, media_ids=["other-file"])
    denied_image = await tool.function(ctx, media_ids=["image-id"])

    assert denied_other == {
        "success": False,
        "error_code": "file_not_found_or_not_accessible",
    }
    assert denied_image == {
        "success": False,
        "error_code": "file_not_found_or_not_accessible",
    }


async def test_read_file_content_truncates_large_content(tmp_path: Path) -> None:
    user_id = uuid4()
    conversation_id = uuid4()
    path = tmp_path / "notes.md"
    path.write_text("abcdef", encoding="utf-8")
    store = FakeMediaAssetStore()
    await store.create(
        asset=MediaAsset(
            media_id="owned-file",
            user_id=user_id,
            conversation_id=conversation_id,
            media_type=MediaAssetType.DOCUMENT,
            storage_key=str(path),
            content_type="text/markdown",
            size_bytes=path.stat().st_size,
            source_channel="telegram",
            metadata={"original_filename": "notes.md"},
        )
    )
    toolsets = build_file_reading_toolsets(
        AppSettings(environment="test", document_max_extracted_chars=3),
        media_asset_store=store,
    )
    tool = cast(Any, toolsets[0]).tools["readFileContent"]
    ctx = cast(
        Any,
        FakeRunContext(deps={"user_id": user_id, "conversation_id": conversation_id}),
    )

    result = await tool.function(ctx, media_ids=["owned-file"])

    assert result["success"] is True
    assert result["data"]["files"][0]["content"] == "abc"
    assert result["data"]["files"][0]["truncated"] is True


async def test_read_file_content_returns_error_for_pdf_without_extractable_text(
    tmp_path: Path,
) -> None:
    user_id = uuid4()
    conversation_id = uuid4()
    path = tmp_path / "blank.pdf"
    _write_pdf(path, text="")
    store = FakeMediaAssetStore()
    await store.create(
        asset=MediaAsset(
            media_id="blank-pdf",
            user_id=user_id,
            conversation_id=conversation_id,
            media_type=MediaAssetType.DOCUMENT,
            storage_key=str(path),
            content_type="application/pdf",
            size_bytes=path.stat().st_size,
            source_channel="telegram",
            metadata={"original_filename": "blank.pdf"},
        )
    )
    toolsets = build_file_reading_toolsets(
        AppSettings(environment="test"),
        media_asset_store=store,
    )
    tool = cast(Any, toolsets[0]).tools["readFileContent"]
    ctx = cast(
        Any,
        FakeRunContext(deps={"user_id": user_id, "conversation_id": conversation_id}),
    )

    result = await tool.function(ctx, media_ids=["blank-pdf"])

    assert result == {
        "success": False,
        "error_code": "empty_file_content",
    }


async def test_read_file_content_returns_error_for_pptx_without_extractable_text(
    tmp_path: Path,
) -> None:
    user_id = uuid4()
    conversation_id = uuid4()
    content_type = "application/vnd.openxmlformats-officedocument.presentationml.presentation"
    path = tmp_path / "blank.pptx"
    _write_pptx(path, include_content=False)
    store = FakeMediaAssetStore()
    await store.create(
        asset=MediaAsset(
            media_id="blank-pptx",
            user_id=user_id,
            conversation_id=conversation_id,
            media_type=MediaAssetType.DOCUMENT,
            storage_key=str(path),
            content_type=content_type,
            size_bytes=path.stat().st_size,
            source_channel="telegram",
            metadata={"original_filename": "blank.pptx"},
        )
    )
    toolsets = build_file_reading_toolsets(
        AppSettings(environment="test"),
        media_asset_store=store,
    )
    tool = cast(Any, toolsets[0]).tools["readFileContent"]
    ctx = cast(
        Any,
        FakeRunContext(deps={"user_id": user_id, "conversation_id": conversation_id}),
    )

    result = await tool.function(ctx, media_ids=["blank-pptx"])

    assert result == {
        "success": False,
        "error_code": "empty_file_content",
    }


def _write_docx(path: Path) -> None:
    from docx import Document

    document = Document()
    document.add_paragraph("Project brief")
    table = document.add_table(rows=2, cols=2)
    table.cell(0, 0).text = "Name"
    table.cell(0, 1).text = "Value"
    table.cell(1, 0).text = "Count"
    table.cell(1, 1).text = "42"
    document.save(str(path))


def _write_pptx(path: Path, *, include_content: bool) -> None:
    from pptx import Presentation
    from pptx.util import Inches

    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])

    if include_content:
        title = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(8), Inches(0.6))
        title.text = "Project review"

        body = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(8), Inches(1))
        text_frame = body.text_frame
        text_frame.text = "Main goals"
        paragraph = text_frame.add_paragraph()
        paragraph.text = "Increase retention"
        paragraph.level = 1

        table_shape = slide.shapes.add_table(
            rows=2,
            cols=2,
            left=Inches(0.5),
            top=Inches(2.4),
            width=Inches(6),
            height=Inches(1),
        )
        table = table_shape.table
        table.cell(0, 0).text = "Metric"
        table.cell(0, 1).text = "Value"
        table.cell(1, 0).text = "Retention"
        table.cell(1, 1).text = "42%"

    presentation.save(str(path))


def _write_pdf(path: Path, *, text: str) -> None:
    escaped = text.replace("\\", "\\\\").replace("(", "\\(").replace(")", "\\)")
    stream = f"BT /F1 12 Tf 72 720 Td ({escaped}) Tj ET" if text else ""
    stream_bytes = stream.encode("latin-1")
    objects = [
        b"<< /Type /Catalog /Pages 2 0 R >>",
        b"<< /Type /Pages /Kids [3 0 R] /Count 1 >>",
        b"<< /Type /Page /Parent 2 0 R /MediaBox [0 0 612 792] "
        b"/Resources << /Font << /F1 4 0 R >> >> /Contents 5 0 R >>",
        b"<< /Type /Font /Subtype /Type1 /BaseFont /Helvetica >>",
        b"<< /Length "
        + str(len(stream_bytes)).encode("ascii")
        + b" >>\nstream\n"
        + stream_bytes
        + b"\nendstream",
    ]
    parts = [b"%PDF-1.4\n"]
    offsets = [0]
    for index, obj in enumerate(objects, start=1):
        offsets.append(sum(len(part) for part in parts))
        parts.append(f"{index} 0 obj\n".encode("ascii"))
        parts.append(obj)
        parts.append(b"\nendobj\n")
    xref_offset = sum(len(part) for part in parts)
    parts.append(f"xref\n0 {len(objects) + 1}\n".encode("ascii"))
    parts.append(b"0000000000 65535 f \n")
    for offset in offsets[1:]:
        parts.append(f"{offset:010d} 00000 n \n".encode("ascii"))
    parts.append(
        (
            f"trailer\n<< /Size {len(objects) + 1} /Root 1 0 R >>\n"
            f"startxref\n{xref_offset}\n%%EOF\n"
        ).encode("ascii")
    )
    path.write_bytes(b"".join(parts))
